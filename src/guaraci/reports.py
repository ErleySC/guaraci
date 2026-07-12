"""reports.py — Geração de relatórios do app web (PDF, Word, Excel, LaTeX,
PowerPoint).

Extraído de app_quimiometria.py (item 18 da auditoria: quebrar o monólito da
UI web por serviço/aba). Estas funções são PURAS quanto a I/O de terminal —
não importam `streamlit`, recebem (pasta, projeto) e devolvem bytes/BytesIO.
A UI (app_quimiometria.py) apenas as chama e envolve o resultado em
`st.cache_data` + `st.download_button`.
"""
from __future__ import annotations

import io
import logging
import os
import time
from typing import Dict, List

import pandas as pd

import guaraci.pipeline as _pq
from guaraci.app_logic import (
    ler_resumo as _ler_resumo,
    listar_figuras as _listar_figuras,
)
from guaraci.config import NOME_TABELAS
# Parsing do resumo_modelo.txt centralizado (item 19): _ex e o dicionario de
# metricas eram duplicados nos 5 geradores; agora vem de resumo_parse.
from guaraci.resumo_parse import extrair_metrica, parse_metricas_modelo

# Fonte unica de versao (mesmo padrao de app_quimiometria.py: pipeline.__version__).
_APP_VERSION = f"v{getattr(_pq, '__version__', '?')}"


def gerar_pdf_relatorio(pasta: str, projeto: Dict,
                          max_figuras: int = 14) -> io.BytesIO:
    """
    Generates a complete PDF report with fpdf2.
    Structure: Cover | Metrics | Figures (2/page) | References.
    Returns BytesIO ready for st.download_button.
    """
    import unicodedata
    from fpdf import FPDF

    def _a(txt: str) -> str:
        """Removes accents for fpdf2 Latin-1 fonts."""
        return unicodedata.normalize("NFKD", str(txt)).encode("ascii", "ignore").decode("ascii")

    # ── Parse resumo_modelo.txt ───────────────────────────────────────
    resumo_raw = _ler_resumo(pasta) or ""

    metricas = parse_metricas_modelo(resumo_raw)

    imgs = _listar_figuras(pasta)[:max_figuras]

    # ── PDF class ──────────────────────────────────────────────────────
    class RelatorioPDF(FPDF):
        def header(self):
            if self.page_no() == 1:
                return
            self.set_font("Helvetica", "B", 8)
            self.set_text_color(100, 100, 100)
            proj_nome = _a(projeto.get("nome", "Plataforma Quimiometrica"))
            self.cell(130, 6, proj_nome[:60], border=0, align="L")
            self.cell(0, 6, f"Generated: {time.strftime('%Y-%m-%d')}",
                      border=0, align="R", new_x="LMARGIN", new_y="NEXT")
            self.set_draw_color(200, 200, 200)
            self.line(15, self.get_y(), 195, self.get_y())
            self.ln(2)
            self.set_text_color(0, 0, 0)

        def footer(self):
            if self.page_no() == 1:
                return
            self.set_y(-12)
            self.set_font("Helvetica", "I", 8)
            self.set_text_color(140, 140, 140)
            self.cell(0, 6, f"Page {self.page_no()} / {{nb}}", align="C")

    pdf = RelatorioPDF(orientation="P", unit="mm", format="A4")
    pdf.alias_nb_pages()
    pdf.set_margins(15, 20, 15)
    pdf.set_auto_page_break(auto=True, margin=20)

    # ── COVER ────────────────────────────────────────────────────────────
    pdf.add_page()
    # Blue top band
    pdf.set_fill_color(30, 80, 140)
    pdf.rect(0, 0, 210, 55, style="F")
    pdf.set_font("Helvetica", "B", 24)
    pdf.set_text_color(255, 255, 255)
    pdf.set_xy(15, 12)
    pdf.cell(0, 12, "Chemometrics Platform", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 13)
    pdf.set_xy(15, 30)
    pdf.cell(0, 8, "PLS-DA  |  PCA  |  OPLS-DA  |  DD-SIMCA", align="C")

    pdf.set_text_color(0, 0, 0)
    pdf.set_xy(15, 65)
    pdf.set_font("Helvetica", "B", 17)
    pdf.multi_cell(180, 9,
                   _a(projeto.get("nome", "Chemometric Analysis Report")),
                   align="C")

    pdf.ln(8)
    pdf.set_font("Helvetica", "", 11)
    pdf.set_fill_color(240, 244, 250)
    campos_capa = [
        ("Author(s)",      projeto.get("autor", "-")),
        ("Institution",    projeto.get("inst", "-")),
        ("Study type",     projeto.get("tipo", "-")),
        ("Date",           time.strftime("%Y-%m-%d")),
        ("Folder",         os.path.basename(pasta)),
    ]
    for label, valor in campos_capa:
        pdf.set_font("Helvetica", "B", 11)
        pdf.cell(55, 7, f"{label}:", border="B", fill=True)
        pdf.set_font("Helvetica", "", 11)
        pdf.cell(125, 7, _a(str(valor))[:80], border="B", fill=True,
                 new_x="LMARGIN", new_y="NEXT")

    obj = projeto.get("objetivo", "")
    if obj:
        pdf.ln(6)
        pdf.set_font("Helvetica", "B", 11)
        pdf.cell(0, 7, "Objective:", new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", "", 10)
        pdf.multi_cell(180, 5, _a(obj))

    # Cover footer
    pdf.set_y(-30)
    pdf.set_fill_color(30, 80, 140)
    pdf.set_text_color(255, 255, 255)
    pdf.set_font("Helvetica", "I", 8)
    pdf.rect(0, pdf.get_y(), 210, 30, style="F")
    pdf.set_xy(15, pdf.get_y() + 5)
    pdf.cell(0, 5,
             f"Generated by: Chemometrics Platform {_APP_VERSION}  |  GEAAp / UFPA  |  PIBIC Project",
             align="C")

    # ── SECTION 1: METRICS ───────────────────────────────────────────────
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 14)
    pdf.set_text_color(30, 80, 140)
    pdf.cell(0, 9, "1. Executive Summary", new_x="LMARGIN", new_y="NEXT")
    pdf.set_draw_color(30, 80, 140)
    pdf.set_line_width(0.5)
    pdf.line(15, pdf.get_y(), 195, pdf.get_y())
    pdf.set_line_width(0.2)
    pdf.ln(4)

    # Metrics table
    pdf.set_text_color(0, 0, 0)
    pdf.set_font("Helvetica", "B", 9)
    pdf.set_fill_color(30, 80, 140)
    pdf.set_text_color(255, 255, 255)
    pdf.cell(100, 7, "Metric", border=1, fill=True)
    pdf.cell(75, 7, "Value (internal CV validation)", border=1, fill=True,
             new_x="LMARGIN", new_y="NEXT")
    pdf.set_text_color(0, 0, 0)
    for i, (k, v) in enumerate(metricas.items()):
        if i % 2 == 0:
            pdf.set_fill_color(240, 245, 255)
        else:
            pdf.set_fill_color(252, 252, 252)
        pdf.set_font("Helvetica", "", 9)
        pdf.cell(100, 6, _a(k), border=1, fill=True)
        pdf.cell(75, 6, _a(str(v)), border=1, fill=True,
                 new_x="LMARGIN", new_y="NEXT")

    # Quality criteria
    pdf.ln(6)
    pdf.set_font("Helvetica", "B", 11)
    pdf.set_text_color(30, 80, 140)
    pdf.cell(0, 8, _a("Quality criteria (PLS-DA — literature):"),
             new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(0, 0, 0)
    criterios = [
        "R2Y > 0.60 and Q2Y > 0.50 → robust predictive model [Eriksson 2008]",
        "Q2Y / R2Y > 0.70 → low overfitting risk",
        "permutation p-value < 0.05 → non-random result (Y-randomization)",
        "Intercept R2Y < 0.40 and intercept Q2Y < 0.05 → geometric validation [Wold 1984]",
        "BCa 95% CI of Bal.Acc does not include 1/n_classes → significant discrimination [Efron 1993]",
    ]
    for c in criterios:
        pdf.set_x(18)
        pdf.cell(4, 5, "-")
        pdf.multi_cell(173, 5, _a(c))

    # ── SECTION 2: FIGURES ────────────────────────────────────────────────
    if imgs:
        # 2 figures per page, stacked vertically
        fig_w = 175
        fig_h = 100   # max height per figure
        cap_h = 6     # caption height

        for par in range(0, len(imgs), 2):
            pdf.add_page()
            if par == 0:
                pdf.set_font("Helvetica", "B", 14)
                pdf.set_text_color(30, 80, 140)
                pdf.cell(0, 9, "2. Figures", new_x="LMARGIN", new_y="NEXT")
                pdf.set_draw_color(30, 80, 140)
                pdf.set_line_width(0.5)
                pdf.line(15, pdf.get_y(), 195, pdf.get_y())
                pdf.set_line_width(0.2)
                pdf.ln(3)

            # Top figure
            y_top = pdf.get_y()
            nome1 = os.path.splitext(os.path.basename(imgs[par]))[0]
            try:
                pdf.image(imgs[par], x=15, y=y_top, w=fig_w, h=fig_h)
            except OSError as _e_img:   # arquivo ausente/corrompido/formato invalido
                logging.getLogger(__name__).warning(
                    "PDF: figura nao embutida (%s): %s", imgs[par], _e_img)
            pdf.set_xy(15, y_top + fig_h + 1)
            pdf.set_font("Helvetica", "I", 7)
            pdf.set_text_color(80, 80, 80)
            pdf.cell(fig_w, cap_h,
                     f"Fig. {par + 1}: {_a(nome1[:70])}",
                     align="C", new_x="LMARGIN", new_y="NEXT")

            # Bottom figure (if it exists)
            if par + 1 < len(imgs):
                pdf.ln(3)
                y_bot = pdf.get_y()
                nome2 = os.path.splitext(os.path.basename(imgs[par + 1]))[0]
                try:
                    pdf.image(imgs[par + 1], x=15, y=y_bot, w=fig_w, h=fig_h)
                except OSError as _e_img:
                    logging.getLogger(__name__).warning(
                        "PDF: figura nao embutida (%s): %s",
                        imgs[par + 1], _e_img)
                pdf.set_xy(15, y_bot + fig_h + 1)
                pdf.set_font("Helvetica", "I", 7)
                pdf.cell(fig_w, cap_h,
                         f"Fig. {par + 2}: {_a(nome2[:70])}",
                         align="C")

            pdf.set_text_color(0, 0, 0)

    # ── SECTION 3: REFERENCES ────────────────────────────────────────────
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 14)
    pdf.set_text_color(30, 80, 140)
    pdf.cell(0, 9, "3. References", new_x="LMARGIN", new_y="NEXT")
    pdf.set_draw_color(30, 80, 140)
    pdf.set_line_width(0.5)
    pdf.line(15, pdf.get_y(), 195, pdf.get_y())
    pdf.set_line_width(0.2)
    pdf.ln(5)
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(0, 0, 0)

    refs = [
        "Bylesjo, M. et al. (2006) OPLS discriminant analysis: combining the strengths of PLS-DA and SIMCA. J. Chemom. 20:341-351.",
        "Chong, I.G. & Jun, C.H. (2005) Performance of some variable selection methods when multicollinearity is present. Chemom. Intell. Lab. Syst. 78:103-112.",
        "Efron, B. & Tibshirani, R.J. (1993) An Introduction to the Bootstrap. CRC Press, Boca Raton.",
        "Eriksson, L. et al. (2008) CV-ANOVA for significance testing of PLS and OPLS models. J. Chemom. 22:594-600.",
        "Pomerantsev, A.L. & Rodionova, O.Y. (2020) Concept and role of extreme objects in PCA/SIMCA. J. Chemom. 34:e3227.",
        "Rajalahti, T. et al. (2009) Discriminating variable test and selectivity ratio plot: quantitative tools for interpretation and variable (biomarker) selection in complex spectral or chromatographic profiles. Anal. Chem. 81:2581-2590.",
        "Rinnan, A. et al. (2009) Review of the most common pre-processing techniques for near-infrared spectra. TrAC Trends Anal. Chem. 28:1201-1222.",
        "Tracy, N.D. et al. (1992) Multivariate control charts for individual observations. J. Qual. Tech. 24:88-95.",
        "Wold, S. (1978) Cross-validatory estimation of the number of components in factor and principal components models. Technometrics 20:397-405.",
    ]
    for i, ref in enumerate(refs, 1):
        pdf.multi_cell(0, 5, f"[{i}] {_a(ref)}")
        pdf.ln(1)

    pdf.ln(8)
    pdf.set_font("Helvetica", "I", 8)
    pdf.set_text_color(120, 120, 120)
    pdf.multi_cell(0, 5, _a(
                   f"Report automatically generated by the Chemometrics Platform {_APP_VERSION}. "
                   "Engine: pipeline.py | Interface: app_quimiometria.py. "
                   "GEAAp/UFPA — PIBIC Project."))

    buf = io.BytesIO()
    buf.write(pdf.output())
    buf.seek(0)
    return buf


def gerar_word_relatorio(pasta: str, projeto: Dict,
                           max_figuras: int = 14) -> io.BytesIO:
    """
    Generates an editable Word report (.docx) with python-docx.
    Same structure as the PDF: cover, metrics, figures, references.
    """
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    resumo_raw = _ler_resumo(pasta) or ""

    metricas = parse_metricas_modelo(resumo_raw)
    imgs = _listar_figuras(pasta)[:max_figuras]

    doc = Document()

    # ── Heading styles ──
    for i, (level, size) in enumerate([(0, 20), (1, 16), (2, 13)]):
        style = doc.styles[f"Heading {i + 1}"]
        style.font.size = Pt(size)           # type: ignore[union-attr]
        style.font.color.rgb = RGBColor(30, 80, 140)  # type: ignore[union-attr]
        style.font.bold = True               # type: ignore[union-attr]

    # ── COVER ──────────────────────────────────────────────────────────────
    doc.add_heading("Chemometrics Platform", 0)
    doc.add_heading(projeto.get("nome", "Chemometric Analysis Report"), 1)

    t_capa = doc.add_table(rows=5, cols=2)
    t_capa.style = "Table Grid"
    campos_capa = [
        ("Author(s)",      projeto.get("autor", "-")),
        ("Institution",    projeto.get("inst", "GEAAp / UFPA")),
        ("Study type",     projeto.get("tipo", "-")),
        ("Date",           time.strftime("%Y-%m-%d %H:%M")),
        ("Folder",         os.path.basename(pasta)),
    ]
    for i, (label, valor) in enumerate(campos_capa):
        c0 = t_capa.cell(i, 0)
        c1 = t_capa.cell(i, 1)
        c0.text = label
        c1.text = str(valor)
        c0.paragraphs[0].runs[0].bold = True

    obj = projeto.get("objetivo", "")
    if obj:
        doc.add_heading("Objective", 2)
        doc.add_paragraph(obj)

    doc.add_page_break()

    # ── SECTION 1: METRICS ─────────────────────────────────────────────────
    doc.add_heading("1. Executive Summary — Model Metrics", 1)

    t_met = doc.add_table(rows=1 + len(metricas), cols=2)
    t_met.style = "Table Grid"
    hdr = t_met.rows[0].cells
    hdr[0].text = "Metric"
    hdr[1].text = "Value (internal CV validation)"
    for cell in hdr:
        run = cell.paragraphs[0].runs
        if run:
            run[0].bold = True
        cell.paragraphs[0].paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER

    for i, (k, v) in enumerate(metricas.items(), 1):
        t_met.cell(i, 0).text = k
        t_met.cell(i, 1).text = str(v)

    doc.add_heading("Quality criteria", 2)
    criterios = [
        "R2Y > 0.60 and Q2Y > 0.50: robust predictive model [Eriksson 2008]",
        "Q2Y / R2Y > 0.70: low overfitting risk",
        "permutation p-value < 0.05: non-random result (Y-randomization)",
        "Intercept R2Y < 0.40 and intercept Q2Y < 0.05 [Wold 1984]",
        "BCa 95% CI of Bal.Acc does not include 1/n_classes [Efron 1993]",
    ]
    for c in criterios:
        doc.add_paragraph(c, style="List Bullet")

    # ── SECTION 2: FIGURES ──────────────────────────────────────────────────
    if imgs:
        doc.add_page_break()
        doc.add_heading("2. Figures", 1)
        for i, img in enumerate(imgs, 1):
            nome_fig = os.path.splitext(os.path.basename(img))[0]
            try:
                doc.add_picture(img, width=Inches(5.5))
                p = doc.paragraphs[-1]
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                leg = doc.add_paragraph(f"Fig. {i}: {nome_fig}")
                leg.alignment = WD_ALIGN_PARAGRAPH.CENTER
                if leg.runs:
                    leg.runs[0].italic = True
                    leg.runs[0].font.size = Pt(9)
            except OSError:   # arquivo ausente/corrompido/formato invalido
                doc.add_paragraph(f"[Fig. {i}: {nome_fig} — image not available]")
            if i % 2 == 0 and i < len(imgs):
                doc.add_page_break()

    # ── SECTION 3: REFERENCES ─────────────────────────────────────────────
    doc.add_page_break()
    doc.add_heading("3. References", 1)
    refs = [
        "Bylesjo, M. et al. (2006) OPLS discriminant analysis. J. Chemom. 20:341-351.",
        "Chong, I.G. & Jun, C.H. (2005) VIP scores. Chemom. Intell. Lab. Syst. 78:103-112.",
        "Efron, B. & Tibshirani, R.J. (1993) An Introduction to the Bootstrap. CRC Press.",
        "Eriksson, L. et al. (2008) CV-ANOVA for PLS/OPLS significance. J. Chemom. 22:594-600.",
        "Pomerantsev, A.L. & Rodionova, O.Y. (2020) Cooman's Plot. J. Chemom. 34:e3227.",
        "Rajalahti, T. et al. (2009) Selectivity Ratio. Anal. Chem. 81:2581-2590.",
        "Rinnan, A. et al. (2009) NIR pre-processing review. TrAC 28:1201-1222.",
        "Tracy, N.D. et al. (1992) Multivariate control charts. J. Qual. Tech. 24:88-95.",
        "Wold, S. (1978) Cross-validatory estimation. Technometrics 20:397-405.",
    ]
    for i, ref in enumerate(refs, 1):
        p = doc.add_paragraph(f"[{i}] {ref}")
        p.paragraph_format.space_after = Pt(4)

    doc.add_paragraph()
    nota = doc.add_paragraph(
        f"Report generated by the Chemometrics Platform {_APP_VERSION}. "
        "Engine: pipeline.py | GEAAp/UFPA — PIBIC Project."
    )
    if nota.runs:
        nota.runs[0].italic = True
        nota.runs[0].font.size = Pt(8)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def gerar_excel_relatorio(pasta: str) -> io.BytesIO:
    """
    Generates an Excel report with 4 sheets via openpyxl:
      - Metrics: metrics extracted from the summary
      - Identifiers: samples with T2, Q, class (pipeline CSV)
      - VIP_Selection: VIP/SR scores (pipeline CSV, if present)
      - Raw_Summary: full text of resumo_modelo.txt
    """
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    resumo_raw = _ler_resumo(pasta) or ""

    def _ex(padrao: str, default: str = "-") -> str:
        return extrair_metrica(resumo_raw, padrao, default)

    metricas_dict = {
        "Balanced Accuracy (CV)":  _ex(r"[Bb]alanced[_ ]?[Aa]ccuracy.*?[:=]\s*([\d.]+)"),
        "AUC macro OvR":           _ex(r"ROC AUC macro.*?[:=]\s*([\d.]+)"),
        "R2Y":                     _ex(r"\bR2Y\b.*?[:=]\s*([\d.]+)"),
        "Q2Y":                     _ex(r"\bQ2\b.*?[:=]\s*([\d.E+-]+)"),
        "R2X":                     _ex(r"\bR2X\b.*?[:=]\s*([\d.]+)"),
        "Optimal LVs":             _ex(r"LVs?\s+otim[ao].*?[:=]\s*(\d+)"),
        "p-value (permutation)":   _ex(r"p.?value.*?[:=]\s*([\d.E+-]+)"),
        "Preprocessing":           _ex(r"[Pp]re.?[Pp]rocessamento.*?[:=]\s*([A-Za-z0-9_+]+)"),
        "Hotelling T2 UCL (95%)":  _ex(r"[Hh]otelling.*?[:=]\s*([\d.]+)"),
        "Q-residual UCL (95%)":    _ex(r"Q.residual.*?[:=]\s*([\d.E+-]+)"),
        "n samples (training)":    _ex(r"[Nn]\s+treino.*?[:=]\s*(\d+)"),
        "n classes":               _ex(r"[Nn]\.?\s*[Cc]lasses.*?[:=]\s*(\d+)"),
    }

    _AZUL      = PatternFill("solid", fgColor="1E508C")
    _AZUL_CLAR = PatternFill("solid", fgColor="DCE6F5")
    _BRANCO    = PatternFill("solid", fgColor="FFFFFF")
    _HDR_FONT  = Font(bold=True, color="FFFFFF", size=10)
    _BODY_FONT = Font(size=9)
    _THIN      = Side(style="thin", color="AAAAAA")
    _BORDER    = Border(left=_THIN, right=_THIN, top=_THIN, bottom=_THIN)

    def _cabecalho(ws, colunas: List[str]):
        for j, col in enumerate(colunas, 1):
            cell = ws.cell(1, j, col)
            cell.font = _HDR_FONT
            cell.fill = _AZUL
            cell.alignment = Alignment(horizontal="center")
            cell.border = _BORDER

    def _preencher_df(ws, df: pd.DataFrame):
        _cabecalho(ws, list(df.columns))
        for i, row in enumerate(df.itertuples(index=False), 2):
            fill = _AZUL_CLAR if i % 2 == 0 else _BRANCO
            for j, val in enumerate(row, 1):
                c = ws.cell(i, j, val)
                c.font = _BODY_FONT
                c.fill = fill
                c.border = _BORDER

    def _preencher_df_ws(ws, df: pd.DataFrame, row_start: int = 1):
        """Variant of _preencher_df with configurable row_start."""
        cols = list(df.columns)
        for j, col in enumerate(cols, 1):
            cell = ws.cell(row_start, j, col)
            cell.font = _HDR_FONT; cell.fill = _AZUL
            cell.alignment = Alignment(horizontal="center"); cell.border = _BORDER
        for i, row in enumerate(df.itertuples(index=False), row_start + 1):
            fill = _AZUL_CLAR if i % 2 == 0 else _BRANCO
            for j, val in enumerate(row, 1):
                c = ws.cell(i, j, val)
                c.font = _BODY_FONT; c.fill = fill; c.border = _BORDER

    def _auto_width(ws):
        for col in ws.columns:
            max_len = max((len(str(c.value)) for c in col if c.value), default=8)
            ws.column_dimensions[col[0].column_letter].width = min(max_len + 3, 40)

    wb = openpyxl.Workbook()

    # ── SHEET 1: Metrics ───────────────────────────────────────────────────
    ws1 = wb.active
    assert ws1 is not None
    ws1.title = "Metrics"
    _cabecalho(ws1, ["Metric", "Value (internal CV)"])
    for i, (k, v) in enumerate(metricas_dict.items(), 2):
        fill = _AZUL_CLAR if i % 2 == 0 else _BRANCO
        for j, val in enumerate([k, v], 1):
            c = ws1.cell(i, j, val)
            c.font = _BODY_FONT
            c.fill = fill
            c.border = _BORDER
    _auto_width(ws1)

    # ── SHEET 2: Identifiers ────────────────────────────────────────────────
    ws2 = wb.create_sheet("Identifiers")
    id_csv = os.path.join(pasta, NOME_TABELAS, "amostras_identificadores.csv")
    if os.path.exists(id_csv):
        try:
            df_id = pd.read_csv(id_csv, sep=";", decimal=",")
            _preencher_df(ws2, df_id)
            _auto_width(ws2)
        except (pd.errors.ParserError, OSError, UnicodeDecodeError):
            ws2.cell(1, 1, "Error reading amostras_identificadores.csv")
    else:
        ws2.cell(1, 1, "File not found (run pipeline Step 5 first).")

    # ── SHEET 3: VIP_Selection ────────────────────────────────────────────────
    ws3 = wb.create_sheet("VIP_Selection")
    vip_csv = os.path.join(pasta, NOME_TABELAS, "etapa4_selecao_variaveis.csv")
    if os.path.exists(vip_csv):
        try:
            df_vip = pd.read_csv(vip_csv, sep=";", decimal=",")
            _preencher_df(ws3, df_vip)
            _auto_width(ws3)
        except (pd.errors.ParserError, OSError, UnicodeDecodeError):
            ws3.cell(1, 1, "Error reading etapa4_selecao_variaveis.csv")
    else:
        ws3.cell(1, 1, "Step 4 (variable selection) was not executed.")

    # ── SHEET 4: Raw_Summary ───────────────────────────────────────────────
    ws4 = wb.create_sheet("Raw_Summary")
    ws4.cell(1, 1, "resumo_modelo.txt — full content").font = Font(bold=True)
    for i, linha in enumerate(resumo_raw.splitlines(), 2):
        ws4.cell(i, 1, linha).font = Font(name="Courier New", size=9)
    ws4.column_dimensions["A"].width = 80

    # ── SHEET 5: Benchmark (if it exists) ─────────────────────────────────
    bench_csv = os.path.join(pasta, NOME_TABELAS, "benchmark_classificadores.csv")
    mc_csv    = os.path.join(pasta, NOME_TABELAS, "monte_carlo_cv.csv")
    if os.path.exists(bench_csv) or os.path.exists(mc_csv):
        ws5 = wb.create_sheet("Benchmark")
        row_cursor = 1
        if os.path.exists(bench_csv):
            try:
                df_bench = pd.read_csv(bench_csv, sep=";", decimal=",")
                ws5.cell(row_cursor, 1,
                         "Auto-Benchmark — Bal.Acc by classifier (GroupKFold)").font = Font(bold=True)
                row_cursor += 1
                _preencher_df_ws(ws5, df_bench, row_start=row_cursor)
                row_cursor += len(df_bench) + 3
                _auto_width(ws5)
            except (pd.errors.ParserError, OSError, UnicodeDecodeError):
                ws5.cell(row_cursor, 1, "Error reading benchmark_classificadores.csv")
                row_cursor += 2
        if os.path.exists(mc_csv):
            try:
                df_mc = pd.read_csv(mc_csv, sep=";", decimal=",")
                ws5.cell(row_cursor, 1,
                         "Monte Carlo CV — 95% CI by percentile").font = Font(bold=True)
                row_cursor += 1
                _preencher_df_ws(ws5, df_mc, row_start=row_cursor)
                _auto_width(ws5)
            except (pd.errors.ParserError, OSError, UnicodeDecodeError):
                ws5.cell(row_cursor, 1, "Error reading monte_carlo_cv.csv")

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf


def gerar_latex_template(pasta: str, projeto: Dict) -> bytes:
    """
    Generates a LaTeX template ready for journals (Talanta, Food Chemistry,
    Journal of Chemometrics). Includes auto-filled metrics, \\includegraphics
    blocks for figures, and complete bibliography.
    Returns UTF-8 bytes (.tex file).
    """

    resumo_raw = _ler_resumo(pasta) or ""

    def _ex(padrao: str, default: str = "-") -> str:
        return extrair_metrica(resumo_raw, padrao, default)

    def _esc(txt: str) -> str:
        """Escapes LaTeX special characters."""
        for old, new in [
            ("\\", r"\textbackslash{}"), ("&", r"\&"), ("%", r"\%"),
            ("$", r"\$"), ("#", r"\#"), ("_", r"\_"),
            ("{", r"\{"), ("}", r"\}"),
            ("~", r"\textasciitilde{}"), ("^", r"\textasciicircum{}"),
        ]:
            txt = txt.replace(old, new)
        return txt

    met = {
        "bal_acc":   _ex(r"[Bb]alanced.*?[:=]\s*([\d.]+)"),
        "auc":       _ex(r"ROC AUC macro.*?[:=]\s*([\d.]+)"),
        "r2y":       _ex(r"\bR2Y\b.*?[:=]\s*([\d.]+)"),
        "q2y":       _ex(r"\bQ2\b.*?[:=]\s*([\d.E+-]+)"),
        "r2x":       _ex(r"\bR2X\b.*?[:=]\s*([\d.]+)"),
        "lvs":       _ex(r"LVs?\s+otim[ao].*?[:=]\s*(\d+)"),
        "perm_p":    _ex(r"p.?value.*?[:=]\s*([\d.E+-]+)"),
        "preproc":   _ex(r"[Pp]re.?[Pp]rocessamento.*?[:=]\s*([A-Za-z0-9_+]+)"),
        "n_train":   _ex(r"[Nn]\s+treino.*?[:=]\s*(\d+)"),
        "n_classes": _ex(r"[Nn]\.?\s*[Cc]lasses.*?[:=]\s*(\d+)"),
    }

    imgs = _listar_figuras(pasta)[:8]
    nome_proj = _esc(projeto.get("nome", "Chemometric Analysis by FT-NIR"))
    autor     = _esc(projeto.get("autor", "Surname, N."))
    inst      = _esc(projeto.get("inst", "GEAAp, Federal University of Para"))

    # Metrics table
    linhas_met = [
        r"        \textbf{Metric} & \textbf{Value (internal CV)} \\",
        r"        \midrule",
        f"        Balanced Accuracy & {met['bal_acc']} \\\\",
        f"        AUC macro (OvR)   & {met['auc']} \\\\",
        f"        $R^2Y$            & {met['r2y']} \\\\",
        f"        $Q^2Y$            & {met['q2y']} \\\\",
        f"        $R^2X$            & {met['r2x']} \\\\",
        f"        Optimal LVs       & {met['lvs']} \\\\",
        f"        $p$ (permutation) & {met['perm_p']} \\\\",
        f"        Preprocessing     & {_esc(met['preproc'])} \\\\",
        f"        $n$ training      & {met['n_train']} \\\\",
        f"        $n$ classes       & {met['n_classes']} \\\\",
    ]
    tabela = "\n".join(linhas_met)

    # Figure blocks
    blocos_fig = []
    for i, img in enumerate(imgs, 1):
        nome_f = os.path.splitext(os.path.basename(img))[0]
        label  = nome_f.replace(" ", "_").replace("-", "_")
        # LaTeX prefers forward slashes
        img_tex = img.replace("\\", "/")
        blocos_fig.append(f"""
\\begin{{figure}}[htbp]
    \\centering
    \\includegraphics[width=0.85\\linewidth]{{{img_tex}}}
    \\caption{{{_esc(nome_f)}. % TODO: add chemical interpretation.}}
    \\label{{fig:{label}}}
\\end{{figure}}""")
    figs_block = "\n".join(blocos_fig)

    tex = f"""% ================================================================
% LaTeX Template — Chemometrics Platform {_APP_VERSION}
% Compatible with: Talanta, Food Chemistry, J. Chemometrics (Elsevier)
% Generated: {time.strftime('%Y-%m-%d %H:%M')}
% ================================================================

\\documentclass[12pt,a4paper]{{article}}

%% Packages
\\usepackage[utf8]{{inputenc}}
\\usepackage[T1]{{fontenc}}
\\usepackage[english]{{babel}}
\\usepackage{{amsmath, amssymb}}
\\usepackage{{graphicx}}
\\usepackage{{booktabs}}
\\usepackage{{longtable}}
\\usepackage[colorlinks=true,citecolor=blue,linkcolor=blue]{{hyperref}}
\\usepackage{{geometry}}
\\usepackage{{caption}}
\\usepackage{{float}}
\\usepackage{{siunitx}}
\\usepackage{{natbib}}

\\geometry{{a4paper, left=2.5cm, right=2.5cm, top=3cm, bottom=3cm}}
\\captionsetup{{font=small, labelfont=bf}}
\\setlength{{\\parindent}}{{0.5cm}}

%% Metadata
\\title{{{nome_proj}}}
\\author{{{autor} \\\\
         \\small {inst}}}
\\date{{\\today}}

\\begin{{document}}
\\maketitle

%% ── Abstract ─────────────────────────────────────────────────────────────
\\begin{{abstract}}
The authenticity of Amazonian vegetable oils was investigated by FT-NIR
spectroscopy combined with chemometric methods (PLS-DA, PCA, OPLS-DA and DD-SIMCA).
The PLS-DA model with {_esc(met['lvs'])} latent variables and {_esc(met['preproc'])}
preprocessing achieved a balanced accuracy of {met['bal_acc']}
($R^2Y = {met['r2y']}$; $Q^2Y = {met['q2y']}$; $p = {met['perm_p']}$,
Y-randomization), evidencing significant discriminatory power.
% TODO: limit to 250 words for Talanta / 200 for Food Chemistry.
\\end{{abstract}}

\\textbf{{Keywords:}} FT-NIR; PLS-DA; chemometrics; authentication; vegetable oils.

%% ── Introduction ──────────────────────────────────────────────────────────
\\section{{Introduction}}
% TODO: contextualize, justify the method and state the objective in ~3 paragraphs.
Adulteration of vegetable oils is a global food safety problem \\citep{{rinnan2009}}.
FT-NIR spectroscopy combined with chemometrics has shown potential as a rapid,
non-destructive, low-cost technique for oil authentication
\\citep{{bylesjo2006, eriksson2008}}.

%% ── Material and Methods ──────────────────────────────────────────────────
\\section{{Material and Methods}}

\\subsection{{Samples and spectral acquisition}}
% TODO: describe number of samples, instrument (ABB MB3600 or similar),
% spectral range, resolution, number of scans, temperature.
FT-NIR spectra were acquired in the range \\SIrange{{4000}}{{10000}}{{\\per\\centi\\meter}}.

\\subsection{{Spectral preprocessing}}
The {_esc(met['preproc'])} preprocessing was selected based on pipeline comparison
(balanced accuracy in CV validation) according to \\citet{{rinnan2009}}.

\\subsection{{PLS-DA modelling}}
The PLS-DA model was calibrated with {met['lvs']} latent variables, selected
by group-aware cross-validation (GroupKFold, grouping technical replicates to
prevent data leakage) \\citep{{chong2005}}.

\\subsection{{Statistical validation}}
Robustness was assessed by:
\\begin{{itemize}}
    \\item \\textbf{{Y-randomization}} (200 permutations) \\citep{{eriksson2008}};
    \\item \\textbf{{Wold's test}} (intercepts $R^2Y < 0.40$ and $Q^2Y < 0.05$);
    \\item \\textbf{{CV-ANOVA}} by Eriksson \\citep{{eriksson2008}};
    \\item \\textbf{{BCa Bootstrap}} 95\\% for balanced accuracy \\citep{{efron1993}}.
\\end{{itemize}}

%% ── Results and Discussion ──────────────────────────────────────────────
\\section{{Results and Discussion}}

\\subsection{{Model performance}}
Performance metrics are presented in Table~\\ref{{tab:metricas}}.

\\begin{{table}}[htbp]
    \\centering
    \\caption{{PLS-DA model performance metrics.}}
    \\label{{tab:metricas}}
    \\begin{{tabular}}{{ll}}
        \\toprule
{tabela}
        \\bottomrule
    \\end{{tabular}}
\\end{{table}}

\\subsection{{Figures}}
% Figures generated by the pipeline. Insert chemical interpretation in each caption.

{figs_block}

%% ── Conclusion ───────────────────────────────────────────────────────────
\\section{{Conclusion}}
% TODO: summarise results and implications in 1-2 paragraphs.
The PLS-DA model with {_esc(met['preproc'])} preprocessing achieved satisfactory
performance ($Q^2Y = {met['q2y']}$; $p = {met['perm_p']}$), demonstrating the
viability of FT-NIR spectroscopy for rapid authentication of Amazonian vegetable oils.

%% ── Acknowledgements ──────────────────────────────────────────────────────
\\section*{{Acknowledgements}}
% TODO: CNPq, CAPES, PIBIC/UFPA, laboratory.
To GEAAp/UFPA and CNPq for financial support (PIBIC Project).

%% ── Referencias ─────────────────────────────────────────────────────────
\\bibliographystyle{{elsarticle-num}}  %% Elsevier (Talanta, Food Chemistry)
%% \\bibliographystyle{{apalike}}       %% Alternativa APA

\\begin{{thebibliography}}{{99}}

\\bibitem{{bylesjo2006}}
Bylesjo,~M. et~al. (2006).
OPLS discriminant analysis: combining the strengths of PLS-DA and SIMCA.
\\textit{{Journal of Chemometrics}}, 20(8--10), 341--351.

\\bibitem{{chong2005}}
Chong,~I.G. \\& Jun,~C.H. (2005).
Performance of some variable selection methods when multicollinearity is present.
\\textit{{Chemometrics and Intelligent Laboratory Systems}}, 78(1--2), 103--112.

\\bibitem{{efron1993}}
Efron,~B. \\& Tibshirani,~R.J. (1993).
\\textit{{An Introduction to the Bootstrap}}. CRC Press.

\\bibitem{{eriksson2008}}
Eriksson,~L., Trygg,~J. \\& Wold,~S. (2008).
CV-ANOVA for significance testing of PLS and OPLS models.
\\textit{{Journal of Chemometrics}}, 22(11--12), 594--600.

\\bibitem{{pomerantsev2020}}
Pomerantsev,~A.L. \\& Rodionova,~O.Y. (2020).
Concept and role of extreme objects in PCA/SIMCA.
\\textit{{Journal of Chemometrics}}, 34(3), e3227.

\\bibitem{{rajalahti2009}}
Rajalahti,~T. et~al. (2009).
Discriminating variable test and selectivity ratio plot.
\\textit{{Analytical Chemistry}}, 81(7), 2581--2590.

\\bibitem{{rinnan2009}}
Rinnan,~\\AA., van~den~Berg,~F. \\& Engelsen,~S.B. (2009).
Review of the most common pre-processing techniques for near-infrared spectra.
\\textit{{TrAC Trends in Analytical Chemistry}}, 28(10), 1201--1222.

\\bibitem{{tracy1992}}
Tracy,~N.D., Young,~J.C. \\& Mason,~R.L. (1992).
Multivariate control charts for individual observations.
\\textit{{Journal of Quality Technology}}, 24(2), 88--95.

\\bibitem{{wold1978}}
Wold,~S. (1978).
Cross-validatory estimation of the number of components.
\\textit{{Technometrics}}, 20(4), 397--405.

\\end{{thebibliography}}

\\end{{document}}
"""
    return tex.encode("utf-8")


# ──────────────────────────────────────────────────────────────────────────
# Execution with live feedback (background thread + progress bar)
# ──────────────────────────────────────────────────────────────────────────


def gerar_pptx_relatorio(pasta: str, projeto: Dict,
                           max_figuras: int = 12) -> io.BytesIO:
    """
    Generates a PowerPoint presentation with professional scientific design.
    Palette: Navy #0F172A + Accent #0369A1 (UI Pro Max — B2B Professional).
    Structure: Cover | Methodology | Metrics | Figures | Benchmark | Conclusions
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # ── UI Pro Max — B2B Professional palette ──────────────────────────────
    _NAVY   = RGBColor(0x0F, 0x17, 0x2A)
    _SLATE  = RGBColor(0x33, 0x41, 0x55)
    _ACCENT = RGBColor(0x03, 0x69, 0xA1)
    _LIGHT  = RGBColor(0xF8, 0xFA, 0xFC)
    _WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    _MUTED  = RGBColor(0x64, 0x74, 0x8B)

    resumo_raw = _ler_resumo(pasta) or ""

    def _ex(padrao: str, default: str = "—") -> str:
        return extrair_metrica(resumo_raw, padrao, default)

    metricas = {
        "Balanced Accuracy (CV)": _ex(r"[Bb]alanced[_ ]?[Aa]cc.*?[:=]\s*([\d.]+)"),
        "AUC macro OvR":          _ex(r"ROC AUC macro.*?[:=]\s*([\d.]+)"),
        "R2Y":                    _ex(r"\bR2Y\b.*?[:=]\s*([\d.]+)"),
        "Q2Y":                    _ex(r"\bQ2\b.*?[:=]\s*([\d.E+-]+)"),
        "Optimal LVs":            _ex(r"LVs?\s+otim.*?[:=]\s*(\d+)"),
        "Preprocessing":          _ex(r"[Pp]re.?[Pp]rocess.*?[:=]\s*([A-Za-z0-9_+]+)"),
        "N samples":              _ex(r"[Nn]\s+treino.*?[:=]\s*(\d+)"),
        "N classes":              _ex(r"[Nn]\.?\s*[Cc]lasses.*?[:=]\s*(\d+)"),
        "p permutation":          _ex(r"p.?value.*?[:=]\s*([\d.E+-]+)"),
    }

    # ── Layout helpers ─────────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W = prs.slide_width   # Emu — assigned above, never None
    H = prs.slide_height  # Emu — assigned above, never None
    assert W is not None and H is not None  # satisfies Pylance Optional[Emu] stub

    def _rect(slide, l, t, w, h, fill_rgb: RGBColor):
        from pptx.util import Emu
        shape = slide.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
        shape.fill.solid(); shape.fill.fore_color.rgb = fill_rgb
        shape.line.fill.background()
        return shape

    def _txt(slide, texto: str, l, t, w, h,
             bold=False, size=18, color=_WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
        tb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
        tf = tb.text_frame; tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run(); run.text = texto
        run.font.bold = bold; run.font.size = Pt(size)
        run.font.color.rgb = color
        return tb

    def _slide_novo(layout_idx=6):
        layout = prs.slide_layouts[layout_idx]
        return prs.slides.add_slide(layout)

    def _fundo(slide, cor=_LIGHT):
        bg = slide.background; fill = bg.fill
        fill.solid(); fill.fore_color.rgb = cor

    def _barra_topo(slide, titulo: str):
        _rect(slide, 0, 0, int(W), int(Inches(1.0)), _NAVY)
        _rect(slide, 0, int(Inches(1.0)), int(W), int(Inches(0.06)), _ACCENT)
        _txt(slide, titulo, int(Inches(0.4)), int(Inches(0.15)),
             int(W - Inches(0.8)), int(Inches(0.75)),
             bold=True, size=22, color=_WHITE)

    def _rodape(slide):
        _rect(slide, 0, int(H - Inches(0.4)), int(W), int(Inches(0.4)), _SLATE)
        data_str = time.strftime("%Y-%m-%d")
        inst = projeto.get("inst", "GEAAp / UFPA")
        _txt(slide,
             f"{inst}  •  Chemometrics Platform  •  {data_str}",
             int(Inches(0.3)), int(H - Inches(0.35)),
             int(W - Inches(0.6)), int(Inches(0.35)),
             size=9, color=_LIGHT, align=PP_ALIGN.CENTER)

    # ── SLIDE 1: Cover ─────────────────────────────────────────────────────
    slide1 = _slide_novo()
    _fundo(slide1, _NAVY)
    # Horizontal accent band
    _rect(slide1, 0, int(Inches(3.5)), int(W), int(Inches(0.12)), _ACCENT)
    # Main title
    _txt(slide1, projeto.get("nome", "Chemometrics Platform"),
         int(Inches(1.0)), int(Inches(1.2)),
         int(W - Inches(2.0)), int(Inches(1.5)),
         bold=True, size=36, color=_WHITE)
    # Subtitle
    _txt(slide1, projeto.get("tipo", "Chemometric Analysis FT-NIR"),
         int(Inches(1.0)), int(Inches(2.8)),
         int(W - Inches(2.0)), int(Inches(0.8)),
         size=20, color=RGBColor(0xCB, 0xD5, 0xE1))
    # Metadata
    autor  = projeto.get("autor", "")
    inst   = projeto.get("inst", "GEAAp / UFPA")
    data_s = time.strftime("%Y-%m-%d")
    _txt(slide1, f"{autor}\n{inst}\n{data_s}",
         int(Inches(1.0)), int(Inches(4.0)),
         int(Inches(6.0)), int(Inches(2.0)),
         size=14, color=RGBColor(0x94, 0xA3, 0xB8))

    # ── SLIDE 2: Methodology ──────────────────────────────────────────────
    slide2 = _slide_novo()
    _fundo(slide2)
    _barra_topo(slide2, "Methodology")
    _rodape(slide2)

    metod_items = [
        f"Preprocessing: {metricas['Preprocessing']}",
        f"Main classifier: PLS-DA ({metricas['Optimal LVs']} optimal LVs)",
        "Cross-validation: GroupKFold anti-leakage of replicates (mae_id)",
        "Statistical tests: Y-permutation, Wold (R2Y/Q2Y), CV-ANOVA",
        "Variable selection: iPLS, VIP >= 1, SR top-20%, sPLS-DA",
        "External validation: stratified holdout (pure samples always in training)",
        "Benchmark: PLS-DA vs SVM RBF vs RF vs GBM vs XGBoost",
        f"Data: {metricas['N samples']} samples  |  {metricas['N classes']} classes",
    ]
    for i, item in enumerate(metod_items):
        _txt(slide2, f"• {item}",
             int(Inches(0.7)), int(Inches(1.3) + i * Inches(0.62)),
             int(W - Inches(1.4)), int(Inches(0.6)),
             size=13, color=_SLATE)

    # ── SLIDE 3: Performance Metrics ───────────────────────────────────────
    slide3 = _slide_novo()
    _fundo(slide3)
    _barra_topo(slide3, "Results — Performance Metrics")
    _rodape(slide3)

    met_exibir = [
        ("Balanced Accuracy (CV)", metricas["Balanced Accuracy (CV)"], _ACCENT),
        ("AUC macro OvR",          metricas["AUC macro OvR"],          _NAVY),
        ("R2Y",                    metricas["R2Y"],                    _SLATE),
        ("Q2Y",                    metricas["Q2Y"],                    _SLATE),
        ("Optimal LVs",            metricas["Optimal LVs"],            _MUTED),
        ("p permutation",          metricas["p permutation"],          _MUTED),
    ]
    cols = 3; w_box = int((W - Inches(1.0)) / cols)
    h_box = int(Inches(1.6))
    for idx, (lbl, val, cor) in enumerate(met_exibir):
        col = idx % cols; row = idx // cols
        x = int(Inches(0.5) + col * w_box)
        y = int(Inches(1.3)  + row * Inches(2.0))
        _rect(slide3, x, y, w_box - int(Inches(0.15)), h_box, cor)
        _txt(slide3, val, x + int(Inches(0.2)), y + int(Inches(0.2)),
             w_box - int(Inches(0.4)), int(Inches(0.8)),
             bold=True, size=28, color=_WHITE, align=PP_ALIGN.CENTER)
        _txt(slide3, lbl, x + int(Inches(0.1)), y + int(Inches(0.95)),
             w_box - int(Inches(0.2)), int(Inches(0.55)),
             size=11, color=_LIGHT, align=PP_ALIGN.CENTER)

    # ── SLIDES 4+: Figures ────────────────────────────────────────────────
    imgs = _listar_figuras(pasta)
    # Prioritize relevant figures
    prioridade = ["scores","confus","vip","pca","outlier","splot","cooman",
                  "roc","hca","opls","ddsimca","benchmark","shap","monte_carlo"]
    def _prio(p: str) -> int:
        nome_l = os.path.basename(p).lower()
        for i, t in enumerate(prioridade):
            if t in nome_l: return i
        return 99
    imgs_sorted = sorted(imgs, key=_prio)[:max_figuras]

    for i in range(0, len(imgs_sorted), 2):
        slide_f = _slide_novo()
        _fundo(slide_f)
        _barra_topo(slide_f, "Results — Figures")
        _rodape(slide_f)
        for j, img_path in enumerate(imgs_sorted[i:i+2]):
            x = int(Inches(0.3) + j * Inches(6.4))
            try:
                slide_f.shapes.add_picture(img_path,
                    Emu(x), Emu(int(Inches(1.15))),
                    width=Emu(int(Inches(6.1))),
                    height=Emu(int(Inches(5.6))))
            except OSError:   # arquivo ausente/corrompido/formato invalido
                _txt(slide_f, f"[Figura: {os.path.basename(img_path)}]",
                     x, int(Inches(2.0)), int(Inches(6.0)), int(Inches(1.0)),
                     size=11, color=_MUTED)
            cap = os.path.splitext(os.path.basename(img_path))[0].replace("_", " ")
            _txt(slide_f, cap, x, int(H - Inches(0.85)),
                 int(Inches(6.0)), int(Inches(0.45)),
                 size=9, color=_MUTED, align=PP_ALIGN.CENTER)

    # ── SLIDE Benchmark (if CSV exists) ──────────────────────────────────
    bench_csv = os.path.join(pasta, NOME_TABELAS, "benchmark_classificadores.csv")
    if os.path.exists(bench_csv):
        try:
            df_b = pd.read_csv(bench_csv, sep=";", decimal=",")
            slide_b = _slide_novo()
            _fundo(slide_b)
            _barra_topo(slide_b, "Classifier Benchmark")
            _rodape(slide_b)
            # Table
            cols_b = list(df_b.columns)
            n_rows = min(len(df_b) + 1, 8)
            col_w  = int((W - Inches(1.0)) / len(cols_b))
            row_h  = int(Inches(0.52))
            # Header
            for ci, col_n in enumerate(cols_b):
                _rect(slide_b, int(Inches(0.5) + ci * col_w),
                      int(Inches(1.3)), col_w - 2, row_h, _NAVY)
                _txt(slide_b, str(col_n),
                     int(Inches(0.5) + ci * col_w + Inches(0.05)),
                     int(Inches(1.35)), col_w - int(Inches(0.1)),
                     int(Inches(0.45)), bold=True, size=10, color=_WHITE)
            # Data rows
            fill_alts = [_LIGHT, RGBColor(0xE8, 0xEC, 0xF1)]
            for ri, row_data in enumerate(df_b.itertuples(index=False)):
                if ri >= 7: break
                fy = int(Inches(1.3) + (ri + 1) * row_h)
                for ci, val in enumerate(row_data):
                    _rect(slide_b, int(Inches(0.5) + ci * col_w),
                          fy, col_w - 2, row_h, fill_alts[ri % 2])
                    _txt(slide_b, str(val),
                         int(Inches(0.5) + ci * col_w + Inches(0.05)),
                         fy + int(Inches(0.05)),
                         col_w - int(Inches(0.1)), int(Inches(0.45)),
                         size=10, color=_SLATE)
        except Exception as _e_slide:  # noqa: BLE001 -- slide opcional
            # (leitura CSV + N formas desenhadas); erro logado, o restante
            # do PPTX (capa, figuras, conclusoes) continua sendo gerado.
            logging.getLogger(__name__).warning(
                "PPTX: slide de benchmark nao gerado: %s", _e_slide)

    # ── Final SLIDE: Conclusions ───────────────────────────────────────────
    slide_c = _slide_novo()
    _fundo(slide_c, _NAVY)
    _rect(slide_c, 0, int(Inches(3.2)), int(W), int(Inches(0.1)), _ACCENT)
    _txt(slide_c, "Conclusions",
         int(Inches(1.0)), int(Inches(1.0)),
         int(W - Inches(2.0)), int(Inches(0.9)),
         bold=True, size=30, color=_WHITE)
    ba   = metricas.get("Balanced Accuracy (CV)", "—")
    auc  = metricas.get("AUC macro OvR", "—")
    r2y  = metricas.get("R2Y", "—"); q2y = metricas.get("Q2Y", "—")
    conc = [
        f"• MSC-SG-MC pipeline achieved Balanced Accuracy = {ba} (GroupKFold anti-leakage)",
        f"• AUC macro OvR = {auc}  |  R2Y = {r2y}  |  Q2Y = {q2y}",
        f"• {metricas['N samples']} samples, {metricas['N classes']} classes "
          f"of Amazonian vegetable oils — FT-NIR spectroscopy",
        f"• Validation: Y-permutation (p = {metricas['p permutation']}), "
          f"Wold, CV-ANOVA and external holdout",
        "• Platform exports .joblib models for prediction of new samples",
    ]
    for i, c in enumerate(conc):
        _txt(slide_c, c,
             int(Inches(1.0)), int(Inches(2.0) + i * Inches(0.78)),
             int(W - Inches(2.0)), int(Inches(0.75)),
             size=14, color=RGBColor(0xCB, 0xD5, 0xE1))

    _rodape(slide_c)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
